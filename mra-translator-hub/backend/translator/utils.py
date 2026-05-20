from __future__ import annotations

import uuid
from pathlib import Path

from django.conf import settings
from django.core.mail import send_mail
from django.core.files.uploadedfile import UploadedFile
from docx import Document
from gtts import gTTS
from PIL import Image
from pptx import Presentation
from PyPDF2 import PdfReader

import pytesseract

try:
    from deep_translator import GoogleTranslator
except ImportError:
    GoogleTranslator = None

LANGUAGE_ALIASES = {
    "auto": "auto",
    "zh": "chinese (simplified)",
    "zh-cn": "chinese (simplified)",
    "zh-tw": "chinese (traditional)",
    "he": "hebrew",
}

GTTS_LANGUAGES = {
    "zh": "zh-CN",
    "zh-cn": "zh-CN",
    "zh-CN": "zh-CN",
    "zh-tw": "zh-TW",
    "zh-TW": "zh-TW",
    "jw": "jv",
    "he": "iw",
    "auto": "en",
}


def normalize_lang(code: str) -> str:
    value = (code or "auto").strip()
    return LANGUAGE_ALIASES.get(value.lower(), value)


def chunk_text(text: str, max_chars: int = 4200) -> list[str]:
    paragraphs = text.splitlines() or [text]
    chunks: list[str] = []
    current = ""
    for paragraph in paragraphs:
        candidate = f"{current}\n{paragraph}".strip()
        if len(candidate) <= max_chars:
            current = candidate
        else:
            if current:
                chunks.append(current)
            current = paragraph[:max_chars]
    if current:
        chunks.append(current)
    return chunks


def translate_text(text: str, source_lang: str = "auto", target_lang: str = "ta") -> str:
    if not text or not text.strip():
        raise ValueError("source_text is required")
    if GoogleTranslator is None:
        raise RuntimeError("deep-translator is not installed. Run: py -m pip install -r requirements.txt")

    translator = GoogleTranslator(source=normalize_lang(source_lang), target=normalize_lang(target_lang))
    translated_chunks = [translator.translate(chunk) for chunk in chunk_text(text)]
    return "\n".join(translated_chunks)


def extract_text(uploaded_file: UploadedFile) -> str:
    extension = Path(uploaded_file.name).suffix.lower()
    uploaded_file.seek(0)

    if extension == ".txt":
        return uploaded_file.read().decode("utf-8", errors="ignore")

    if extension == ".pdf":
        reader = PdfReader(uploaded_file)
        return "\n".join(page.extract_text() or "" for page in reader.pages).strip()

    if extension == ".docx":
        document = Document(uploaded_file)
        return "\n".join(paragraph.text for paragraph in document.paragraphs).strip()

    if extension == ".pptx":
        presentation = Presentation(uploaded_file)
        text_runs: list[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs).strip()

    if extension in {".jpg", ".jpeg", ".png"}:
        if settings.TESSERACT_CMD:
            pytesseract.pytesseract.tesseract_cmd = settings.TESSERACT_CMD
        image = Image.open(uploaded_file)
        return pytesseract.image_to_string(image).strip()

    raise ValueError("Unsupported file type. Use TXT, PDF, DOCX, PPTX, JPG, JPEG, or PNG.")


def save_translated_file(original_name: str, translated_text: str) -> tuple[str, str]:
    output_dir = settings.MEDIA_ROOT / "translated"
    output_dir.mkdir(parents=True, exist_ok=True)
    safe_stem = Path(original_name).stem.replace(" ", "_")[:80] or "file"
    output_name = f"translated_{safe_stem}_{uuid.uuid4().hex[:8]}.txt"
    output_path = output_dir / output_name
    output_path.write_text(translated_text, encoding="utf-8")
    return output_name, f"{settings.MEDIA_URL}translated/{output_name}"


def create_tts_mp3(text: str, lang: str = "ta") -> tuple[str, str]:
    if not text or not text.strip():
        raise ValueError("text is required")
    output_dir = settings.MEDIA_ROOT / "tts"
    output_dir.mkdir(parents=True, exist_ok=True)
    file_name = f"speech_{uuid.uuid4().hex[:10]}.mp3"
    output_path = output_dir / file_name
    tts_lang = GTTS_LANGUAGES.get(lang, GTTS_LANGUAGES.get(lang.lower(), lang))
    audio = gTTS(text=text[:4500], lang=tts_lang)
    audio.save(str(output_path))
    return file_name, f"{settings.MEDIA_URL}tts/{file_name}"


def send_admin_email(file_name: str, source_lang: str, target_lang: str, status: str, user_email: str = "") -> None:
    subject = f"New File Translated: {file_name}"
    body = (
        f"User uploaded {file_name}.\n"
        f"Source: {source_lang} -> Target: {target_lang}.\n"
        f"Status: {status}.\n"
        f"User email: {user_email or 'Not provided'}."
    )
    send_mail(subject, body, settings.DEFAULT_FROM_EMAIL, [settings.ADMIN_EMAIL], fail_silently=True)